"""
Generate an 8-slide hackathon MVP deck for the Zagreb prompt-to-map outing app.
Uses python-pptx only (vector shapes + text).

Run: uv run python build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# Dark theme (night / map)
C_BG = RGBColor(0x0D, 0x11, 0x17)
C_PANEL = RGBColor(0x16, 0x1B, 0x22)
C_ACCENT = RGBColor(0xF7, 0x81, 0x66)
C_MUTED = RGBColor(0x8B, 0x94, 0x9E)
C_TEXT = RGBColor(0xF0, 0xF6, 0xFC)
C_MAP = RGBColor(0x23, 0x88, 0x36)
C_PIN = RGBColor(0xDA, 0x36, 0x33)
C_BLOCK = RGBColor(0x30, 0x36, 0x3D)


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _fill_line(shape, color: RGBColor, width_pt: float = 1.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _set_text(
    shape,
    text: str,
    *,
    size_pt: float = 14,
    bold: bool = False,
    color: RGBColor = C_TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_bg(slide, prs: Presentation) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _fill(r, C_BG)


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    rng = __import__("random").Random(42)
    for _ in range(24):
        x = Inches(0.2 + rng.random() * 12.5)
        y = Inches(0.4 + rng.random() * 3.0)
        w = Inches(0.12 + rng.random() * 0.35)
        h = Inches(0.15 + rng.random() * 0.45)
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        _fill(b, C_BLOCK)
        b.line.color.rgb = C_MUTED
        b.line.width = Pt(0.5)

    for dx in (-0.35, 0, 0.45):
        pin = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0 + dx * 2), Inches(2.0), Inches(0.22), Inches(0.22))
        _fill(pin, C_PIN)
        pin.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        pin.line.width = Pt(1)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.1))
    _set_text(tb, "PromptMap Zagreb", size_pt=44, bold=True, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.45), Inches(11.7), Inches(0.65))
    _set_text(
        tb,
        "Natural language → places that match, on one map",
        size_pt=18,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.5))
    _set_text(tb, "your vibe → pins on the map", size_pt=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.55))
    _set_text(tb, "Hackathon MVP", size_pt=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.75))
    _set_text(
        tb,
        "Young people in Zagreb: no single place for everything",
        size_pt=22,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.55))
    _set_text(
        tb,
        "Events, bars, pubs, clubs — info is scattered across many sites and channels",
        size_pt=13,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )

    labels = ["Events\nsites", "Social\nmedia", "Forums", "Bar lists", "Word of\nmouth"]
    cx, cy = Inches(6.5), Inches(4.0)
    r_circle = Inches(2.0)
    import math

    for i, lab in enumerate(labels):
        ang = math.pi * 0.85 - (math.pi * 0.7) * i / (len(labels) - 1)
        ox = cx + r_circle * math.cos(ang)
        oy = cy + r_circle * math.sin(ang)
        el = slide.shapes.add_shape(MSO_SHAPE.OVAL, ox - Inches(0.42), oy - Inches(0.42), Inches(0.84), Inches(0.84))
        _fill(el, C_PANEL)
        el.line.color.rgb = C_MUTED
        el.line.width = Pt(2)
        tbx = slide.shapes.add_textbox(ox - Inches(0.4), oy - Inches(0.35), Inches(0.8), Inches(0.7))
        _set_text(tbx, lab, size_pt=11, align=PP_ALIGN.CENTER)

    q = slide.shapes.add_textbox(cx - Inches(0.35), cy - Inches(0.35), Inches(0.7), Inches(0.7))
    _set_text(q, "?", size_pt=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.65))
    _set_text(tb, "One site. One prompt. One map.", size_pt=22, bold=True, align=PP_ALIGN.CENTER)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.2), Inches(10.9), Inches(0.95)
    )
    _fill(bar, C_PANEL)
    bar.line.color.rgb = C_ACCENT
    bar.line.width = Pt(2.5)
    tbx = slide.shapes.add_textbox(Inches(1.45), Inches(2.35), Inches(10.4), Inches(0.65))
    _set_text(
        tbx,
        "Type what you want to do (any language you use)…",
        size_pt=14,
        color=C_MUTED,
    )

    arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.35), Inches(3.35), Inches(0.65), Inches(0.75))
    _fill_line(arr, C_ACCENT, 0)
    arr.line.fill.background()

    tbx = slide.shapes.add_textbox(Inches(5.0), Inches(3.05), Inches(3.5), Inches(0.35))
    _set_text(tbx, "match & rank", size_pt=11, color=C_ACCENT, align=PP_ALIGN.CENTER)

    m = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.45), Inches(11.3), Inches(1.05))
    _fill(m, RGBColor(0x0D, 0x28, 0x18))
    m.line.color.rgb = C_MAP
    m.line.width = Pt(2)
    tbx = slide.shapes.add_textbox(Inches(1.2), Inches(4.65), Inches(10.9), Inches(0.65))
    _set_text(
        tbx,
        "Interactive map: pins that fit the vibe / request",
        size_pt=15,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.65))
    _set_text(tb, "How it works (concept)", size_pt=20, bold=True, align=PP_ALIGN.CENTER)

    labels = [("User\nprompt", 0.9), ("Understand\nintent & style", 3.15), ("Places &\nevents data", 5.4), ("Map\nresults", 7.65)]
    y = Inches(2.85)
    w, h = Inches(1.85), Inches(1.25)
    boxes = []
    for text, x_in in labels:
        x = Inches(x_in)
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        _fill(r, C_PANEL)
        r.line.color.rgb = C_ACCENT
        r.line.width = Pt(1.8)
        boxes.append((x + w / 2, y + h / 2))
        tbx = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.15), w - Inches(0.16), h - Inches(0.3))
        _set_text(tbx, text, size_pt=12, bold=True, align=PP_ALIGN.CENTER)

    mid_y = y + h / 2
    for i in range(3):
        x1 = boxes[i][0] + w / 2
        x2 = boxes[i + 1][0] - w / 2
        c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, mid_y, x2, mid_y)
        c.line.color.rgb = C_MUTED
        c.line.width = Pt(2.5)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.55))
    _set_text(
        tb,
        "MVP: prove the loop end-to-end for one city (Zagreb)",
        size_pt=12,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )


def slide_example(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.55))
    _set_text(tb, "Example", size_pt=20, bold=True, align=PP_ALIGN.CENTER)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.55))
    _set_text(tb, '"Želim ići u alternativni klub"', size_pt=16, color=C_ACCENT, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.45))
    _set_text(
        tb,
        '→ "I want to go to an alternative club"',
        size_pt=13,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )

    m = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(2.15), Inches(11.6), Inches(4.0))
    _fill(m, RGBColor(0x0E, 0x1F, 0x14))
    m.line.color.rgb = C_MAP
    m.line.width = Pt(2)

    rng = __import__("random").Random(7)
    for _ in range(45):
        px = Inches(1.0 + rng.random() * 11.2)
        py = Inches(2.35 + rng.random() * 3.5)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px, py, Inches(0.04), Inches(0.04))
        _fill(dot, RGBColor(0x1F, 0x3D, 0x2A))
        dot.line.fill.background()

    for px, py in [(2.3, 3.0), (3.9, 4.0), (5.5, 3.3), (7.2, 4.1), (4.4, 5.0)]:
        pin = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(py), Inches(0.28), Inches(0.28))
        _fill(pin, C_PIN)
        pin.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        pin.line.width = Pt(1)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.55))
    _set_text(
        tb,
        "Alternative clubs as map pins — pick one and go",
        size_pt=12,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )


def slide_mvp(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.6))
    _set_text(tb, "MVP (this hackathon)", size_pt=20, bold=True, align=PP_ALIGN.CENTER)

    cards = [
        ("Prompt UI", "Natural-language input", 0.8, 3.35),
        ("Matching", "Vibe / intent → venues", 6.5, 3.35),
        ("Map view", "See options geographically", 0.8, 5.15),
        ("Zagreb", "One city, real need", 6.5, 5.15),
    ]
    for title, sub, x, y in cards:
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.9), Inches(1.35))
        _fill(r, C_PANEL)
        r.line.color.rgb = C_MAP
        r.line.width = Pt(1.5)
        tbx = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(5.5), Inches(0.55))
        _set_text(tbx, title, size_pt=15, bold=True, align=PP_ALIGN.CENTER)
        tbx = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.72), Inches(5.5), Inches(0.5))
        _set_text(tbx, sub, size_pt=11, color=C_MUTED, align=PP_ALIGN.CENTER)


def slide_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.9))
    _set_text(tb, "Live demo", size_pt=32, bold=True, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.65), Inches(11.7), Inches(0.65))
    _set_text(
        tb,
        "We will walk through: prompt → results on the map",
        size_pt=15,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.65), Inches(11.7), Inches(0.55))
    _set_text(
        tb,
        "1. Open the app   →   2. Enter a prompt   →   3. Explore pins",
        size_pt=14,
        bold=True,
        color=C_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.65))
    _set_text(tb, "Questions?", size_pt=24, align=PP_ALIGN.CENTER)


def slide_thanks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.0))
    _set_text(tb, "Thank you", size_pt=40, bold=True, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.75), Inches(11.7), Inches(0.65))
    _set_text(
        tb,
        "PromptMap Zagreb — find your night out, in your own words",
        size_pt=14,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.5))
    _set_text(tb, "Team · Hackathon MVP", size_pt=13, color=C_ACCENT, align=PP_ALIGN.CENTER)


def build(output_path: Path | None = None) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_flow(prs)
    slide_example(prs)
    slide_mvp(prs)
    slide_demo(prs)
    slide_thanks(prs)

    out = output_path or Path(__file__).resolve().parent / "Zagreb_PromptMap_Hackathon_MVP.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote: {path}")
